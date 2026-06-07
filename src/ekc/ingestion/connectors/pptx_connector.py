"""
PPTX connector using python-pptx.
Extracts slide title + body text + speaker notes.
"""
import logging
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
from src.ekc.ingestion.connectors.base import SourceConnector, RawDocument
from src.ekc.db.models import SourceType

logger = logging.getLogger(__name__)


class PPTXConnector(SourceConnector):

    def can_handle(self, path: str) -> bool:
        return path.lower().endswith(".pptx")

    def extract(self, path: str) -> list[RawDocument]:
        logger.info(f"Extracting PPTX: {path}")
        try:
            prs = Presentation(path)
            mtime = datetime.fromtimestamp(os.path.getmtime(path))
            base_title = os.path.splitext(os.path.basename(path))[0]
            docs = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                parts = []
                slide_title = ""

                # Slide title
                if slide.shapes.title and slide.shapes.title.text.strip():
                    slide_title = slide.shapes.title.text.strip()
                    parts.append(f"Slide {slide_num}: {slide_title}")

                # Body text from all shapes
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    if shape == slide.shapes.title:
                        continue
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)

                # Speaker notes (if substantive)
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes and len(notes) > 30:
                        parts.append(f"Notes: {notes}")

                content = "\n".join(parts)
                if not content.strip():
                    continue

                docs.append(RawDocument(
                    title=f"{base_title} — Slide {slide_num}: {slide_title}",
                    content=content,
                    source_type=SourceType.pptx,
                    source_path=path,
                    last_modified=mtime,
                    metadata={
                        "slide_number": slide_num,
                        "section_title": slide_title,
                        "heading_path": f"Slide {slide_num}",
                    },
                ))

            logger.info(f"  extracted {len(docs)} slides")
            return docs

        except Exception as e:
            logger.error(f"PPTX extraction failed for {path}: {e}")
            return []